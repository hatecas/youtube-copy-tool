import os
import re
import uuid
import json
import asyncio
import threading
import tempfile
import subprocess
import shutil
from datetime import datetime
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pathlib import Path
from dotenv import load_dotenv

# .env 파일 찾기: backend/.env 또는 프로젝트 루트/.env
_env_path = Path(__file__).parent / ".env"
if not _env_path.exists():
    _env_path = Path(__file__).parent.parent / ".env"
if _env_path.exists():
    load_dotenv(_env_path)
    print(f"[설정] .env 로드 완료: {_env_path}")
else:
    load_dotenv()
    print("[설정] .env 파일을 찾을 수 없습니다. 환경변수에서 직접 읽습니다.")

app = Flask(__name__)
CORS(app, origins=["*"])

# 제작 파일 저장 디렉토리
ASSETS_DIR = Path(__file__).parent / "assets"
ASSETS_DIR.mkdir(exist_ok=True)

# ============================================================
# 외부 서비스 클라이언트
# ============================================================

# Claude AI
anthropic_client = None
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
if ANTHROPIC_API_KEY:
    try:
        import anthropic
        anthropic_client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
        print(f"[설정] Claude AI 연동: 활성 (키: {ANTHROPIC_API_KEY[:20]}...)")
    except ImportError:
        print("[설정] Claude AI 연동: 실패 - 'pip install anthropic' 필요")
else:
    print("[설정] Claude AI 연동: 비활성 - ANTHROPIC_API_KEY가 .env에 없습니다")

# Supabase
supabase_client = None
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY") or os.getenv("SUPABASE_KEY")
if SUPABASE_URL and SUPABASE_KEY:
    try:
        from supabase import create_client
        supabase_client = create_client(SUPABASE_URL, SUPABASE_KEY)
        print("[설정] Supabase 연동: 활성")
    except ImportError:
        print("[설정] Supabase 연동: 실패 - 'pip install supabase' 필요")
    except Exception as e:
        print(f"[설정] Supabase 연동: 실패 - {e}")
        print("[설정] → Supabase 없이 인메모리 모드로 실행합니다")
else:
    print("[설정] Supabase 연동: 비활성 (인메모리 모드)")


# ============================================================
# DB 헬퍼 (Supabase 또는 인메모리 폴백)
# ============================================================

_memory_store = {}


def db_save_project(project: dict):
    """프로젝트 저장"""
    if supabase_client:
        row = {
            "id": project["id"],
            "status": project["status"],
            "video_urls": project.get("video_urls", []),
            "analyzed_videos": json.dumps(project.get("videos"), ensure_ascii=False) if project.get("videos") else None,
            "topics": json.dumps(project.get("topics"), ensure_ascii=False) if project.get("topics") else None,
            "selected_topic_id": project.get("selected_topic_id"),
            "generated_content": json.dumps(project.get("generated_content"), ensure_ascii=False) if project.get("generated_content") else None,
        }
        supabase_client.table("projects").upsert(row).execute()
    else:
        _memory_store[project["id"]] = project


def db_get_project(project_id: str):
    """프로젝트 조회"""
    if supabase_client:
        res = supabase_client.table("projects").select("*").eq("id", project_id).execute()
        if res.data:
            row = res.data[0]
            return {
                "id": row["id"],
                "status": row["status"],
                "video_urls": row.get("video_urls", []),
                "videos": json.loads(row["analyzed_videos"]) if row.get("analyzed_videos") else [],
                "topics": json.loads(row["topics"]) if row.get("topics") else [],
                "selected_topic_id": row.get("selected_topic_id"),
                "generated_content": json.loads(row["generated_content"]) if row.get("generated_content") else None,
                "confirmed_content": json.loads(row["confirmed_content"]) if row.get("confirmed_content") else None,
                "production_assets": json.loads(row["production_assets"]) if row.get("production_assets") else None,
                "error_message": row.get("error_message"),
                "created_at": row.get("created_at"),
            }
        return None
    else:
        return _memory_store.get(project_id)


def db_update_project(project_id: str, updates: dict):
    """프로젝트 부분 업데이트"""
    if supabase_client:
        row = {}
        if "status" in updates:
            row["status"] = updates["status"]
        if "videos" in updates:
            row["analyzed_videos"] = json.dumps(updates["videos"], ensure_ascii=False)
        if "topics" in updates:
            row["topics"] = json.dumps(updates["topics"], ensure_ascii=False)
        if "selected_topic_id" in updates:
            row["selected_topic_id"] = updates["selected_topic_id"]
        if "generated_content" in updates:
            row["generated_content"] = json.dumps(updates["generated_content"], ensure_ascii=False)
        if "confirmed_content" in updates:
            row["confirmed_content"] = json.dumps(updates["confirmed_content"], ensure_ascii=False)
        if "production_assets" in updates:
            row["production_assets"] = json.dumps(updates["production_assets"], ensure_ascii=False)
        if row:
            supabase_client.table("projects").update(row).eq("id", project_id).execute()
        # error_message 컬럼이 없을 수 있으므로 별도 처리
        if "error_message" in updates:
            try:
                supabase_client.table("projects").update(
                    {"error_message": updates["error_message"]}
                ).eq("id", project_id).execute()
            except Exception:
                print(f"[경고] error_message 저장 실패 - Supabase에 컬럼이 없을 수 있음")
    else:
        project = _memory_store.get(project_id)
        if project:
            project.update(updates)


# ============================================================
# YouTube 유틸리티
# ============================================================

def extract_video_id(url: str):
    """YouTube URL에서 video ID 추출"""
    patterns = [
        r'(?:youtube\.com\/watch\?v=)([a-zA-Z0-9_-]{11})',
        r'(?:youtu\.be\/)([a-zA-Z0-9_-]{11})',
        r'(?:youtube\.com\/embed\/)([a-zA-Z0-9_-]{11})',
        r'(?:youtube\.com\/shorts\/)([a-zA-Z0-9_-]{11})',
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None


def get_video_info(video_id: str) -> dict:
    """YouTube Data API로 영상 정보 가져오기"""
    api_key = os.getenv("YOUTUBE_API_KEY")

    if not api_key:
        return {
            "id": video_id,
            "title": f"영상 ({video_id})",
            "thumbnail": f"https://img.youtube.com/vi/{video_id}/hqdefault.jpg",
            "channelName": "채널명",
            "viewCount": 0,
            "publishedAt": datetime.now().isoformat(),
            "duration": "00:00",
            "description": "",
        }

    try:
        from googleapiclient.discovery import build
        youtube = build("youtube", "v3", developerKey=api_key)

        response = youtube.videos().list(
            part="snippet,statistics,contentDetails",
            id=video_id
        ).execute()

        if not response.get("items"):
            return None

        item = response["items"][0]
        snippet = item["snippet"]
        stats = item.get("statistics", {})
        content = item.get("contentDetails", {})

        return {
            "id": video_id,
            "title": snippet.get("title", ""),
            "thumbnail": f"https://img.youtube.com/vi/{video_id}/hqdefault.jpg",
            "channelName": snippet.get("channelTitle", ""),
            "viewCount": int(stats.get("viewCount", 0)),
            "publishedAt": snippet.get("publishedAt", ""),
            "duration": content.get("duration", ""),
            "description": snippet.get("description", ""),
        }
    except Exception as e:
        print(f"YouTube API 오류: {e}")
        return {
            "id": video_id,
            "title": f"영상 ({video_id})",
            "thumbnail": f"https://img.youtube.com/vi/{video_id}/hqdefault.jpg",
            "channelName": "채널명",
            "viewCount": 0,
            "publishedAt": datetime.now().isoformat(),
            "duration": "00:00",
            "description": "",
        }


def get_transcript(video_id: str) -> str:
    """YouTube 자막/대본 추출 (인트로~5분)"""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi

        ytt_api = YouTubeTranscriptApi()
        transcript_list = ytt_api.fetch(video_id, languages=["ko", "en"])

        # 첫 5분(300초)까지만 추출
        texts = []
        for entry in transcript_list:
            if entry.start > 300:
                break
            texts.append(entry.text)

        return " ".join(texts) if texts else "자막을 찾을 수 없습니다."

    except Exception as e:
        print(f"자막 추출 오류 ({video_id}): {e}")
        return "자막을 추출할 수 없습니다."


def calc_view_ratios(videos: list) -> list[dict]:
    """최근 영상 대비 조회수 비율 계산"""
    if not videos:
        return videos

    sorted_by_date = sorted(
        videos,
        key=lambda v: v.get("publishedAt", ""),
        reverse=True,
    )
    latest_views = sorted_by_date[0]["viewCount"] or 1

    for v in videos:
        v["viewRatio"] = round(v["viewCount"] / latest_views, 2)

    videos.sort(key=lambda x: x["viewCount"], reverse=True)
    return videos


# ============================================================
# AI 분석 (Claude API)
# ============================================================

def analyze_with_ai(videos: list, exclude_topics = None) -> list[dict]:
    """Claude AI로 영상 분석 후 주제 3개 추천"""

    if not anthropic_client:
        return _fallback_analyze(videos)

    video_summaries = []
    for i, v in enumerate(videos, 1):
        transcript_preview = v.get("transcript", "")[:1500]
        video_summaries.append(
            f"[영상 {i}]\n"
            f"- 제목: {v['title']}\n"
            f"- 채널: {v['channelName']}\n"
            f"- 조회수: {v['viewCount']:,}회\n"
            f"- 최근 영상 대비 조회수 비율: {v.get('viewRatio', 1.0)}배\n"
            f"- 업로드: {v['publishedAt']}\n"
            f"- 인트로~5분 대본:\n{transcript_preview}"
        )

    video_block = "\n\n".join(video_summaries)

    exclude_instruction = ""
    if exclude_topics:
        exclude_instruction = (
            f"\n\n중요: 이전에 이미 추천했던 아래 주제들은 제외하고 새로운 주제를 추천해주세요:\n"
            + "\n".join(f"- {t}" for t in exclude_topics)
        )

    prompt = f"""당신은 YouTube 콘텐츠 전략 전문가이자 데이터 분석가입니다.
아래 레퍼런스 영상들을 꼼꼼히 분석하여, 이 영상들을 카피(벤치마킹)해서 만들 수 있는 새 영상 주제를 정확히 3개 추천해주세요.

[분석 기준 - 반드시 모두 반영]
1. 각 영상의 제목에서 공통 키워드, 패턴, 화법(숫자 사용, 질문형, 도발형 등)을 파악
2. 조회수가 높은 영상과 낮은 영상의 차이점 분석 (어떤 요소가 조회수를 끌어올렸는지)
3. 최근 영상 대비 조회수 비율이 높은 영상에 더 가중치 부여
4. 대본(인트로~5분)에서 훅(hook) 방식, 스토리 전개 구조, 시청자 반응 유도 방식을 정밀 분석
5. 채널의 전반적인 콘텐츠 방향성과 시청자 성향 파악
6. 각 주제는 "그대로 촬영에 들어갈 수 있을 정도로" 구체적이어야 함{exclude_instruction}

[레퍼런스 영상 데이터]
{video_block}

아래 JSON 배열 형식으로만 응답해주세요 (다른 텍스트 없이):
[
  {{
    "topic": "구체적인 영상 주제 (실제 영상 제목처럼 구체적으로, 20~40자)",
    "target": "타겟 시청자층 (연령대, 직업/상황, 관심사를 구체적으로 명시)",
    "reasoning": "이 주제를 추천하는 이유를 3~5문장으로 상세하게 작성. 반드시 레퍼런스 영상의 구체적인 데이터(어떤 영상의 조회수가 얼마였는지, 어떤 키워드가 효과적이었는지, 대본에서 어떤 훅이 사용되었는지)를 근거로 포함할 것.",
    "confidence": 0~100 사이의 적중 확률 (레퍼런스 데이터 근거 기반으로 산출)
  }},
  ...
]"""

    try:
        message = anthropic_client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2000,
            messages=[{"role": "user", "content": prompt}],
        )

        response_text = message.content[0].text.strip()

        if response_text.startswith("```"):
            response_text = re.sub(r"^```(?:json)?\s*", "", response_text)
            response_text = re.sub(r"\s*```$", "", response_text)

        topics = json.loads(response_text)

        for t in topics:
            t["id"] = str(uuid.uuid4())

        return topics[:3]

    except Exception as e:
        print(f"Claude AI 분석 오류: {e}")
        return _fallback_analyze(videos)


def generate_with_ai(topic: dict, videos: list) -> dict:
    """Claude AI로 콘텐츠 생성 (썸네일 텍스트, 제목, 대본)"""

    if not anthropic_client:
        return _fallback_generate(topic)

    video_summaries = []
    for i, v in enumerate(videos[:5], 1):
        transcript_preview = v.get("transcript", "")[:1000]
        video_summaries.append(
            f"[영상 {i}] {v['title']} (조회수 {v['viewCount']:,}회)\n"
            f"대본 일부: {transcript_preview}"
        )

    video_block = "\n\n".join(video_summaries)

    prompt = f"""당신은 100만 구독자 YouTube 콘텐츠 크리에이터이자 최고의 카피라이팅 전문가입니다.
아래 주제와 레퍼런스 영상 분석을 바탕으로 실전에서 바로 사용할 수 있는 완성도 높은 콘텐츠를 생성해주세요.

[선택된 주제]
- 주제: {topic['topic']}
- 타겟: {topic['target']}

[레퍼런스 영상 분석]
{video_block}

아래 JSON 형식으로만 응답해주세요 (다른 텍스트 없이):
{{
  "thumbnailTexts": [
    "썸네일 텍스트1",
    "썸네일 텍스트2",
    "썸네일 텍스트3"
  ],
  "titleSuggestions": [
    "제목1",
    "제목2"
  ],
  "scripts": [
    {{
      "targetDescription": "타겟1 구체적 설명",
      "intro": "인트로 대본 전체 (500자 이상)",
      "body": "본문 대본 전체 (3000자 이상, 4개 섹션)"
    }},
    {{
      "targetDescription": "타겟2 구체적 설명",
      "intro": "인트로 대본 전체 (500자 이상)",
      "body": "본문 대본 전체 (3000자 이상, 4개 섹션)"
    }}
  ]
}}

=== 반드시 지켜야 할 품질 기준 ===

[썸네일 텍스트 - 3개]
- 한글 기준 4~8자 (예: "이거 실화?", "월 500만원", "퇴사 각")
- 유튜브에서 가장 클릭률이 높은 패턴 활용: 숫자, 질문, 도발, 공감
- 3개 모두 다른 감정/접근법 사용 (충격형, 공감형, 혜택형)

[제목 - 2개]
- 30~50자 수준으로 충분히 길게
- YouTube SEO 키워드를 자연스럽게 포함
- 검색 노출 + 클릭 유도를 동시에 고려
- 예시 수준: "직장인이 퇴근 후 AI 부업 시작했더니 3개월 만에 월급보다 더 벌게 된 현실적인 방법"
- 반드시 궁금증/감정을 자극하는 요소 포함 (결과 암시, 반전, 구체적 숫자)

[인트로 대본 - 각 500~800자]
- 실제 유튜버가 카메라 앞에서 말하는 듯한 완전한 구어체
- 구조: 강력한 훅(첫 2문장) → 공감/문제 제기 → "이 영상에서는~" 소개 → 시청 유도
- 레퍼런스 영상에서 효과적이었던 훅 방식을 참고하되 원본을 복사하지 말 것
- "~거든요", "~잖아요", "~인데요" 등 실제 말투 사용

[본문 대본 - 각 3000~4000자, 반드시 4개 섹션]
각 섹션은 아래 형식으로 작성:

[섹션 1: 제목]
이 섹션의 전체 대본을 구어체로 작성. 최소 700자 이상.
실제 예시, 구체적인 수치, 비유를 활용하여 설득력 있게 전개.

[섹션 2: 제목]
(동일 형식, 최소 700자)

[섹션 3: 제목]
(동일 형식, 최소 700자)

[섹션 4: 제목]
마무리 섹션. 핵심 요약 + CTA(구독/좋아요 유도) 포함. 최소 700자.

- 각 섹션은 완전한 문장으로 구성된 실제 대본이어야 함 (개요/요약이 아님!)
- "~라고 합니다" "~하면 됩니다" 같은 딱딱한 문체 금지 → "~거든요" "~해보세요" "~있잖아요" 사용
- 두 대본은 톤, 예시, 비유, 깊이가 완전히 달라야 함 (같은 내용을 말투만 바꾸는 것은 안 됨)"""

    try:
        message = anthropic_client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=8000,
            messages=[{"role": "user", "content": prompt}],
        )

        response_text = message.content[0].text.strip()

        if response_text.startswith("```"):
            response_text = re.sub(r"^```(?:json)?\s*", "", response_text)
            response_text = re.sub(r"\s*```$", "", response_text)

        result = json.loads(response_text)

        for script in result.get("scripts", []):
            script["id"] = str(uuid.uuid4())
            script["fullScript"] = f"[인트로]\n{script['intro']}\n\n[본문]\n{script['body']}"

        return result

    except Exception as e:
        print(f"Claude AI 생성 오류: {e}")
        return _fallback_generate(topic)


# ============================================================
# AI 폴백 (API 키 없을 때)
# ============================================================

def _fallback_analyze(videos: list) -> list[dict]:
    """AI 없을 때 기본 주제 추천"""
    titles = ", ".join(v["title"] for v in videos[:3])
    return [
        {
            "id": str(uuid.uuid4()),
            "topic": f"레퍼런스 영상 기반 주제 추천 1",
            "target": "분석 대상 타겟층",
            "reasoning": f"레퍼런스 영상({titles})의 공통 키워드를 분석한 결과입니다. ANTHROPIC_API_KEY를 설정하면 AI 기반 정밀 분석이 제공됩니다.",
            "confidence": 50,
        },
        {
            "id": str(uuid.uuid4()),
            "topic": f"레퍼런스 영상 기반 주제 추천 2",
            "target": "분석 대상 타겟층",
            "reasoning": "AI API 키가 설정되지 않아 기본 추천을 제공합니다.",
            "confidence": 40,
        },
        {
            "id": str(uuid.uuid4()),
            "topic": f"레퍼런스 영상 기반 주제 추천 3",
            "target": "분석 대상 타겟층",
            "reasoning": "AI API 키가 설정되지 않아 기본 추천을 제공합니다.",
            "confidence": 30,
        },
    ]


def _fallback_generate(topic: dict) -> dict:
    """AI 없을 때 기본 콘텐츠 생성"""
    return {
        "thumbnailTexts": [
            "썸네일 텍스트 1 (AI 연결 필요)",
            "썸네일 텍스트 2 (AI 연결 필요)",
            "썸네일 텍스트 3 (AI 연결 필요)",
        ],
        "titleSuggestions": [
            "제목 추천 1 (ANTHROPIC_API_KEY 설정 필요)",
            "제목 추천 2 (ANTHROPIC_API_KEY 설정 필요)",
        ],
        "scripts": [
            {
                "id": str(uuid.uuid4()),
                "targetDescription": topic.get("target", "일반 타겟"),
                "intro": "AI API 키가 설정되면 인트로 대본이 자동 생성됩니다.",
                "body": "[섹션 1: 도입]\nAI 연결 후 자동 생성\n\n[섹션 2: 전개]\nAI 연결 후 자동 생성\n\n[섹션 3: 심화]\nAI 연결 후 자동 생성\n\n[섹션 4: 마무리]\nAI 연결 후 자동 생성",
                "fullScript": "AI API 키가 설정되면 전체 대본이 자동 생성됩니다.",
            },
            {
                "id": str(uuid.uuid4()),
                "targetDescription": "보조 타겟층",
                "intro": "AI API 키가 설정되면 두 번째 인트로 대본이 자동 생성됩니다.",
                "body": "[섹션 1: 도입]\nAI 연결 후 자동 생성\n\n[섹션 2: 전개]\nAI 연결 후 자동 생성\n\n[섹션 3: 심화]\nAI 연결 후 자동 생성\n\n[섹션 4: 마무리]\nAI 연결 후 자동 생성",
                "fullScript": "AI API 키가 설정되면 전체 대본이 자동 생성됩니다.",
            },
        ],
    }


# ============================================================
# 제작 기능 (썸네일, PPT, TTS, 영상)
# ============================================================

def generate_thumbnail_image(text: str, project_id: str) -> str:
    """Pillow로 썸네일 이미지 생성 (1280x720)"""
    from PIL import Image, ImageDraw, ImageFont

    width, height = 1280, 720
    img = Image.new("RGB", (width, height))

    # 그라디언트 배경 생성
    draw = ImageDraw.Draw(img)
    for y in range(height):
        r = int(20 + (y / height) * 30)
        g = int(10 + (y / height) * 15)
        b = int(40 + (y / height) * 60)
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    # 중앙에 강조 원형 그라디언트 효과
    for i in range(200, 0, -1):
        alpha = int(255 * (i / 200) * 0.15)
        x0 = width // 2 - i * 3
        y0 = height // 2 - i * 2
        x1 = width // 2 + i * 3
        y1 = height // 2 + i * 2
        draw.ellipse([x0, y0, x1, y1], fill=(255, 59, 92, alpha) if img.mode == "RGBA" else (min(255, 20 + alpha // 3), 10, min(255, 40 + alpha // 4)))

    # 다시 그리기 (깨끗한 배경)
    img = Image.new("RGB", (width, height))
    draw = ImageDraw.Draw(img)

    # 그라디언트 배경
    for y in range(height):
        ratio = y / height
        r = int(15 + ratio * 25)
        g = int(5 + ratio * 15)
        b = int(35 + ratio * 55)
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    # 하단 악센트 바
    draw.rectangle([0, height - 8, width, height], fill=(255, 59, 92))

    # 폰트 설정 (시스템 폰트 폴백)
    font_size = 80
    font = None
    font_paths = [
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Bold.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ]
    for fp in font_paths:
        if os.path.exists(fp):
            try:
                font = ImageFont.truetype(fp, font_size)
                break
            except Exception:
                continue

    if font is None:
        font = ImageFont.load_default()

    # 텍스트 크기 측정 및 줄바꿈
    lines = []
    current_line = ""
    for char in text:
        test_line = current_line + char
        bbox = draw.textbbox((0, 0), test_line, font=font)
        if bbox[2] - bbox[0] > width - 160:
            lines.append(current_line)
            current_line = char
        else:
            current_line = test_line
    if current_line:
        lines.append(current_line)

    # 텍스트 중앙 배치
    line_height = font_size + 20
    total_height = len(lines) * line_height
    y_start = (height - total_height) // 2

    for i, line in enumerate(lines):
        bbox = draw.textbbox((0, 0), line, font=font)
        text_width = bbox[2] - bbox[0]
        x = (width - text_width) // 2
        y = y_start + i * line_height

        # 텍스트 그림자
        draw.text((x + 3, y + 3), line, fill=(0, 0, 0), font=font)
        # 메인 텍스트
        draw.text((x, y), line, fill=(255, 255, 255), font=font)

    # 저장
    filename = f"thumbnail_{project_id}.png"
    filepath = ASSETS_DIR / filename
    img.save(str(filepath), "PNG")
    print(f"[썸네일] 생성 완료: {filepath}")
    return filename


def generate_ppt(title: str, script: dict, project_id: str) -> str:
    """python-pptx로 PPT 생성"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Emu(12192000)  # 16:9
    prs.slide_height = Emu(6858000)

    def add_bg(slide, r=15, g=15, b=25):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=(240, 240, 245), alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        p.alignment = alignment
        return txBox

    # 1. 표지 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide)
    add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(2),
                 title, font_size=36, bold=True, color=(255, 255, 255),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.8),
                 script.get("targetDescription", ""), font_size=16,
                 color=(136, 136, 160), alignment=PP_ALIGN.CENTER)

    # 2. 인트로 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
                 "INTRO", font_size=14, bold=True, color=(255, 59, 92))
    # 인트로 텍스트를 적절히 나누기
    intro_text = script.get("intro", "")
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(5),
                 intro_text, font_size=16, color=(220, 220, 230))

    # 3. 본문 섹션 슬라이드
    body = script.get("body", "")
    sections = re.split(r'\[섹션\s*\d+[:\s]*([^\]]*)\]', body)

    # sections[0]은 첫 번째 섹션 마커 이전의 텍스트 (보통 빈 문자열)
    # sections[1]은 섹션 제목, sections[2]은 섹션 내용, ...
    section_pairs = []
    i = 1
    while i < len(sections) - 1:
        section_title = sections[i].strip()
        section_content = sections[i + 1].strip()
        if section_title or section_content:
            section_pairs.append((section_title, section_content))
        i += 2

    # 섹션이 파싱되지 않으면 전체 본문을 하나의 슬라이드로
    if not section_pairs:
        section_pairs = [("본문", body)]

    for idx, (sec_title, sec_content) in enumerate(section_pairs, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)

        # 섹션 번호 + 제목
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                     f"SECTION {idx}", font_size=12, bold=True, color=(255, 59, 92))
        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(8.4), Inches(0.8),
                     sec_title if sec_title else f"섹션 {idx}", font_size=24, bold=True,
                     color=(255, 255, 255))

        # 섹션 내용 (글자 수가 많으면 여러 슬라이드로 분할)
        max_chars = 600
        content_chunks = [sec_content[j:j+max_chars] for j in range(0, len(sec_content), max_chars)]

        # 첫 청크는 현재 슬라이드에
        if content_chunks:
            add_text_box(slide, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
                         content_chunks[0], font_size=14, color=(200, 200, 215))

        # 나머지 청크는 추가 슬라이드
        for chunk in content_chunks[1:]:
            extra_slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(extra_slide)
            add_text_box(extra_slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                         f"SECTION {idx} (계속)", font_size=12, bold=True, color=(255, 59, 92))
            add_text_box(extra_slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(5),
                         chunk, font_size=14, color=(200, 200, 215))

    # 4. 아웃트로 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(1),
                 "시청해주셔서 감사합니다", font_size=32, bold=True,
                 color=(255, 255, 255), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.8),
                 "좋아요 & 구독 부탁드립니다!", font_size=18,
                 color=(255, 59, 92), alignment=PP_ALIGN.CENTER)

    # 저장
    filename = f"presentation_{project_id}.pptx"
    filepath = ASSETS_DIR / filename
    prs.save(str(filepath))
    print(f"[PPT] 생성 완료: {filepath}")
    return filename


def generate_tts(script: dict, project_id: str) -> str:
    """Edge TTS로 음성 생성"""
    import edge_tts

    full_text = script.get("fullScript", "")
    if not full_text:
        full_text = f"{script.get('intro', '')}\n\n{script.get('body', '')}"

    # 대본 정리: 섹션 마커 제거, 자연스러운 읽기용으로 변환
    clean_text = re.sub(r'\[섹션\s*\d+[:\s]*[^\]]*\]', '', full_text)
    clean_text = re.sub(r'\[인트로\]|\[본문\]', '', clean_text)
    clean_text = clean_text.strip()

    filename = f"tts_{project_id}.mp3"
    filepath = ASSETS_DIR / filename

    async def _generate():
        communicate = edge_tts.Communicate(clean_text, "ko-KR-SunHiNeural")
        await communicate.save(str(filepath))

    asyncio.run(_generate())
    print(f"[TTS] 생성 완료: {filepath}")
    return filename


def _find_ffmpeg() -> str:
    """사용 가능한 ffmpeg 경로 반환 (시스템 → imageio_ffmpeg 순)"""
    # 1) 시스템 ffmpeg
    if shutil.which("ffmpeg"):
        return "ffmpeg"
    # 2) imageio_ffmpeg 번들
    try:
        import imageio_ffmpeg
        path = imageio_ffmpeg.get_ffmpeg_exe()
        if path:
            return path
    except Exception:
        pass
    raise FileNotFoundError("ffmpeg를 찾을 수 없습니다. 시스템에 ffmpeg를 설치하거나 pip install imageio-ffmpeg를 실행하세요.")


def _find_ffprobe() -> str:
    """사용 가능한 ffprobe 경로 반환"""
    if shutil.which("ffprobe"):
        return "ffprobe"
    # imageio_ffmpeg 번들 ffmpeg와 같은 디렉토리에 ffprobe가 있을 수 있음
    try:
        import imageio_ffmpeg
        ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
        if ffmpeg_path:
            ffprobe_path = ffmpeg_path.replace("ffmpeg", "ffprobe")
            if os.path.exists(ffprobe_path):
                return ffprobe_path
            # ffprobe가 없으면 ffmpeg -i로 대체할 수 있도록 None 반환
    except Exception:
        pass
    return ""


def generate_video(project_id: str, ppt_filename: str, tts_filename: str) -> str:
    """FFmpeg로 PPT 슬라이드 + TTS 음성 → 영상 합성"""
    ppt_path = ASSETS_DIR / ppt_filename
    tts_path = ASSETS_DIR / tts_filename

    # ffmpeg 경로 확인
    ffmpeg_bin = _find_ffmpeg()
    print(f"[영상] ffmpeg 경로: {ffmpeg_bin}")

    # 임시 디렉토리에 슬라이드를 이미지로 변환
    slides_dir = ASSETS_DIR / f"slides_{project_id}"
    slides_dir.mkdir(exist_ok=True)

    try:
        # PPT → 이미지 변환 (LibreOffice 또는 python-pptx + Pillow 폴백)
        slide_images = _ppt_to_images(str(ppt_path), str(slides_dir))

        if not slide_images:
            raise Exception("슬라이드 이미지 변환 실패")

        # TTS 오디오 길이 측정
        audio_duration = _get_audio_duration(str(tts_path))
        if audio_duration <= 0:
            audio_duration = 60  # 기본값

        # 슬라이드당 표시 시간 계산
        slide_duration = audio_duration / len(slide_images)
        slide_duration = max(slide_duration, 3)  # 최소 3초

        print(f"[영상] 슬라이드 {len(slide_images)}개, 오디오 {audio_duration:.1f}초, 슬라이드당 {slide_duration:.1f}초")

        # FFmpeg로 영상 합성
        video_filename = f"video_{project_id}.mp4"
        video_path = ASSETS_DIR / video_filename

        # 슬라이드 리스트 파일 생성
        concat_file = slides_dir / "concat.txt"
        with open(concat_file, "w", encoding="utf-8") as f:
            for img_path in slide_images:
                # Windows 백슬래시를 슬래시로 변환 (FFmpeg 호환)
                safe_path = img_path.replace("\\", "/")
                f.write(f"file '{safe_path}'\n")
                f.write(f"duration {slide_duration}\n")
            # 마지막 이미지 반복 (FFmpeg concat 요구사항)
            safe_last = slide_images[-1].replace("\\", "/")
            f.write(f"file '{safe_last}'\n")

        # FFmpeg 실행
        cmd = [
            ffmpeg_bin, "-y",
            "-f", "concat", "-safe", "0", "-i", str(concat_file),
            "-i", str(tts_path),
            "-vf", "scale=1280:720:force_original_aspect_ratio=decrease,pad=1280:720:(ow-iw)/2:(oh-ih)/2:color=black",
            "-c:v", "libx264", "-preset", "fast", "-crf", "23",
            "-c:a", "aac", "-b:a", "128k",
            "-shortest",
            "-pix_fmt", "yuv420p",
            str(video_path),
        ]

        print(f"[FFmpeg 명령] {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
        if result.returncode != 0:
            print(f"[FFmpeg 오류] returncode={result.returncode}")
            print(f"[FFmpeg stderr] {result.stderr[:1000]}")
            raise Exception(f"FFmpeg 오류: {result.stderr[:200]}")

        # 파일 크기 확인 (0바이트면 실패)
        if not video_path.exists() or video_path.stat().st_size == 0:
            raise Exception("영상 파일이 생성되지 않았거나 크기가 0입니다")

        print(f"[영상] 생성 완료: {video_path} ({video_path.stat().st_size / 1024:.1f}KB)")
        return video_filename

    finally:
        # 임시 슬라이드 이미지 정리
        if slides_dir.exists():
            shutil.rmtree(slides_dir, ignore_errors=True)


def _ppt_to_images(ppt_path: str, output_dir: str) -> list[str]:
    """PPT를 이미지로 변환"""
    # 방법 1: LibreOffice 사용
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "png", "--outdir", output_dir, ppt_path],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0:
            images = sorted(Path(output_dir).glob("*.png"))
            if images:
                return [str(img) for img in images]
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # 방법 2: python-pptx + Pillow로 직접 렌더링
    print("[변환] LibreOffice 없음 - Pillow로 직접 슬라이드 이미지 생성")
    from pptx import Presentation
    from pptx.util import Pt
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(ppt_path)
    images = []

    # 폰트 파일 찾기
    font_file = None
    font_paths = [
        # Windows 한국어 폰트
        "C:/Windows/Fonts/malgun.ttf",      # 맑은 고딕
        "C:/Windows/Fonts/malgunbd.ttf",    # 맑은 고딕 Bold
        "C:/Windows/Fonts/gulim.ttc",       # 굴림
        "C:/Windows/Fonts/batang.ttc",      # 바탕
        "C:/Windows/Fonts/arial.ttf",       # Arial (폴백)
        # Linux 폰트
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        # macOS 폰트
        "/System/Library/Fonts/AppleSDGothicNeo.ttc",
    ]
    for fp in font_paths:
        if os.path.exists(fp):
            font_file = fp
            break

    # Bold 폰트 파일 (Windows)
    bold_font_file = None
    bold_paths = [
        "C:/Windows/Fonts/malgunbd.ttf",
        "C:/Windows/Fonts/arialbd.ttf",
    ]
    for fp in bold_paths:
        if os.path.exists(fp):
            bold_font_file = fp
            break

    def _get_font(size_px, bold=False):
        """크기별 폰트 반환"""
        try:
            path = bold_font_file if (bold and bold_font_file) else font_file
            if path:
                return ImageFont.truetype(path, size_px)
        except Exception:
            pass
        return ImageFont.load_default()

    def _draw_wrapped_text(draw, x, y, text, font, fill, max_width, line_height):
        """줄바꿈 처리하여 텍스트 그리기, 최종 y 위치 반환"""
        current_line = ""
        for char in text:
            test = current_line + char
            bbox = draw.textbbox((0, 0), test, font=font)
            if bbox[2] - bbox[0] > max_width:
                if current_line:
                    draw.text((x, y), current_line, fill=fill, font=font)
                    y += line_height
                current_line = char
            else:
                current_line = test
        if current_line:
            draw.text((x, y), current_line, fill=fill, font=font)
            y += line_height
        return y

    # PPT 슬라이드 크기 (EMU → 비율 계산용)
    slide_w = prs.slide_width or 12192000
    slide_h = prs.slide_height or 6858000
    img_w, img_h = 1280, 720
    scale_x = img_w / slide_w
    scale_y = img_h / slide_h

    for idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (img_w, img_h), (15, 15, 25))
        draw = ImageDraw.Draw(img)

        # 슬라이드 번호 표시 (우측 하단)
        num_font = _get_font(12)
        draw.text((img_w - 60, img_h - 30), f"{idx + 1}/{len(prs.slides)}",
                  fill=(80, 80, 100), font=num_font)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # shape 위치를 이미지 좌표로 변환
            sx = int(shape.left * scale_x) if shape.left else 40
            sy = int(shape.top * scale_y) if shape.top else 40
            sw = int(shape.width * scale_x) if shape.width else (img_w - 80)

            y_pos = sy
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    y_pos += 8
                    continue

                # 폰트 크기/색상/볼드 읽기
                font_size_pt = 18
                is_bold = False
                text_color = (220, 220, 230)

                # paragraph 레벨 폰트 정보
                if paragraph.font and paragraph.font.size:
                    font_size_pt = int(paragraph.font.size / 12700)  # EMU → pt
                if paragraph.font and paragraph.font.bold:
                    is_bold = True
                if paragraph.font and paragraph.font.color and paragraph.font.color.rgb:
                    c = paragraph.font.color.rgb
                    text_color = (c[0], c[1], c[2]) if hasattr(c, '__getitem__') else (220, 220, 230)

                # run 레벨 폰트 정보 (더 정확)
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if run.font.size:
                        font_size_pt = int(run.font.size / 12700)
                    if run.font.bold:
                        is_bold = True
                    if run.font.color and run.font.color.rgb:
                        try:
                            c = run.font.color.rgb
                            text_color = (int(str(c)[0:2], 16), int(str(c)[2:4], 16), int(str(c)[4:6], 16))
                        except Exception:
                            pass

                # pt → 픽셀 변환 (화면 표시용)
                font_size_px = max(int(font_size_pt * 1.5), 14)
                line_height = font_size_px + 6
                used_font = _get_font(font_size_px, is_bold)

                y_pos = _draw_wrapped_text(
                    draw, sx, y_pos, text, used_font, text_color,
                    max_width=sw, line_height=line_height
                )
                y_pos += 4

        img_path = os.path.join(output_dir, f"slide_{idx:03d}.png")
        img.save(img_path)
        images.append(img_path)
        print(f"  [슬라이드 {idx+1}/{len(prs.slides)}] 이미지 생성: {img_path}")

    print(f"[변환] 총 {len(images)}개 슬라이드 이미지 생성 완료")
    return images


def _get_audio_duration(audio_path: str) -> float:
    """오디오 파일의 길이(초) 반환"""
    # 방법 1: ffprobe 사용
    ffprobe_bin = _find_ffprobe()
    if ffprobe_bin:
        try:
            result = subprocess.run(
                [ffprobe_bin, "-v", "quiet", "-show_entries", "format=duration",
                 "-of", "default=noprint_wrappers=1:nokey=1", audio_path],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                return float(result.stdout.strip())
        except Exception:
            pass

    # 방법 2: ffmpeg -i로 duration 추출
    try:
        ffmpeg_bin = _find_ffmpeg()
        result = subprocess.run(
            [ffmpeg_bin, "-i", audio_path, "-f", "null", "-"],
            capture_output=True, text=True, timeout=30
        )
        # stderr에서 Duration: HH:MM:SS.xx 추출
        import re
        match = re.search(r"Duration:\s*(\d+):(\d+):(\d+\.\d+)", result.stderr)
        if match:
            h, m, s = float(match.group(1)), float(match.group(2)), float(match.group(3))
            return h * 3600 + m * 60 + s
    except Exception:
        pass

    # 방법 3: mutagen 라이브러리
    try:
        from mutagen.mp3 import MP3
        audio = MP3(audio_path)
        return audio.info.length
    except Exception:
        pass

    print(f"[경고] 오디오 길이를 알 수 없음: {audio_path}")
    return 0


# ============================================================
# API 엔드포인트
# ============================================================

import re as _re

def _parse_duration_seconds(iso_duration):
    """ISO 8601 duration → 초 변환 (예: PT1M30S → 90)"""
    try:
        h = int((_re.search(r'(\d+)H', iso_duration) or [0,0])[1])
        m = int((_re.search(r'(\d+)M', iso_duration) or [0,0])[1])
        s = int((_re.search(r'(\d+)S', iso_duration) or [0,0])[1])
        return h * 3600 + m * 60 + s
    except Exception:
        return 9999

NEWS_CHANNEL_KEYWORDS = ["뉴스", "news", "MBC", "KBS", "SBS", "YTN", "JTBC", "연합뉴스", "채널A", "TV조선"]

def _is_valid_video(item):
    """쇼츠·뉴스·해외 영상 필터"""
    snippet = item.get("snippet", {})
    content = item.get("contentDetails", {})

    duration_sec = _parse_duration_seconds(content.get("duration", "PT0S"))
    if duration_sec <= 60:
        return False

    channel_title = snippet.get("channelTitle", "")
    if any(kw.lower() in channel_title.lower() for kw in NEWS_CHANNEL_KEYWORDS):
        return False

    lang = snippet.get("defaultAudioLanguage") or snippet.get("defaultLanguage") or ""
    if lang and not lang.startswith("ko"):
        return False

    return True


@app.route("/api/recommend", methods=["GET"])
def recommend():
    """AI 관련 핫한 유튜브 영상 5개 자동 추천"""
    api_key = os.getenv("YOUTUBE_API_KEY")
    if not api_key:
        return jsonify({"error": "YOUTUBE_API_KEY가 설정되지 않았습니다"}), 500

    try:
        from googleapiclient.discovery import build
        youtube = build("youtube", "v3", developerKey=api_key)

        import random
        queries = [
            "AI 유튜브 자동화",
            "AI 콘텐츠 제작",
            "ChatGPT 활용법",
            "AI 부업 만들기",
            "Claude AI 사용법",
            "인공지능 영상 만들기",
            "AI 수익화",
        ]
        query = random.choice(queries)

        response = youtube.search().list(
            part="snippet",
            q=query,
            type="video",
            order="viewCount",
            maxResults=10,
            regionCode="KR",
            relevanceLanguage="ko",
            publishedAfter=datetime.now().strftime("%Y-%m-01T00:00:00Z"),
        ).execute()

        items = response.get("items", [])

        video_ids = [item["id"]["videoId"] for item in items if item["id"].get("videoId")]
        if not video_ids:
            return jsonify({"videos": []})

        detail_response = youtube.videos().list(
            part="snippet,statistics,contentDetails",
            id=",".join(video_ids[:10])
        ).execute()

        videos = []
        for item in detail_response.get("items", []):
            if not _is_valid_video(item):
                continue
            snippet = item["snippet"]
            stats = item.get("statistics", {})
            vid = item["id"]
            videos.append({
                "videoId": vid,
                "title": snippet.get("title", ""),
                "channelName": snippet.get("channelTitle", ""),
                "thumbnail": f"https://img.youtube.com/vi/{vid}/hqdefault.jpg",
                "viewCount": int(stats.get("viewCount", 0)),
                "publishedAt": snippet.get("publishedAt", "")[:10],
                "url": f"https://www.youtube.com/watch?v={vid}",
            })

        videos.sort(key=lambda x: x["viewCount"], reverse=True)
        return jsonify({"videos": videos[:5], "query": query})

    except Exception as e:
        print(f"추천 API 오류: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/recommend-custom", methods=["GET"])
def recommend_custom():
    """키워드 + 채널 기반 맞춤 추천"""
    api_key = os.getenv("YOUTUBE_API_KEY")
    if not api_key:
        return jsonify({"error": "YOUTUBE_API_KEY가 설정되지 않았습니다"}), 500

    keywords = request.args.getlist("keywords")
    channel_ids = request.args.getlist("channelIds")

    if not keywords and not channel_ids:
        return jsonify({"error": "keywords 또는 channelIds 필요"}), 400

    try:
        import random
        from googleapiclient.discovery import build
        youtube = build("youtube", "v3", developerKey=api_key)

        all_videos = []

        if keywords:
            query = random.choice(keywords)
            response = youtube.search().list(
                part="snippet",
                q=query,
                type="video",
                order="viewCount",
                maxResults=10,
                regionCode="KR",
                relevanceLanguage="ko",
                publishedAfter=datetime.now().strftime("%Y-%m-01T00:00:00Z"),
            ).execute()
            video_ids = [item["id"]["videoId"] for item in response.get("items", []) if item["id"].get("videoId")]
            if video_ids:
                detail = youtube.videos().list(part="snippet,statistics,contentDetails", id=",".join(video_ids[:10])).execute()
                for item in detail.get("items", []):
                    if not _is_valid_video(item):
                        continue
                    snippet = item["snippet"]
                    stats = item.get("statistics", {})
                    vid = item["id"]
                    all_videos.append({
                        "videoId": vid,
                        "title": snippet.get("title", ""),
                        "channelName": snippet.get("channelTitle", ""),
                        "thumbnail": f"https://img.youtube.com/vi/{vid}/hqdefault.jpg",
                        "viewCount": int(stats.get("viewCount", 0)),
                        "publishedAt": snippet.get("publishedAt", "")[:10],
                        "url": f"https://www.youtube.com/watch?v={vid}",
                        "source": "keyword",
                        "sourceLabel": query,
                    })

        if channel_ids:
            channel_tags = []

            for channel_id in channel_ids[:5]:
                response = youtube.search().list(
                    part="snippet",
                    channelId=channel_id,
                    type="video",
                    order="date",
                    maxResults=3,
                ).execute()
                video_ids = [item["id"]["videoId"] for item in response.get("items", []) if item["id"].get("videoId")]
                if video_ids:
                    detail = youtube.videos().list(part="snippet,statistics,contentDetails", id=",".join(video_ids)).execute()
                    for item in detail.get("items", []):
                        if not _is_valid_video(item):
                            continue
                        snippet = item["snippet"]
                        stats = item.get("statistics", {})
                        vid = item["id"]
                        all_videos.append({
                            "videoId": vid,
                            "title": snippet.get("title", ""),
                            "channelName": snippet.get("channelTitle", ""),
                            "thumbnail": f"https://img.youtube.com/vi/{vid}/hqdefault.jpg",
                            "viewCount": int(stats.get("viewCount", 0)),
                            "publishedAt": snippet.get("publishedAt", "")[:10],
                            "url": f"https://www.youtube.com/watch?v={vid}",
                            "source": "channel",
                            "sourceLabel": snippet.get("channelTitle", ""),
                        })
                        tags = snippet.get("tags", [])
                        channel_tags.extend(tags[:5])

            if channel_tags:
                import random
                similar_query = random.choice(channel_tags)
                sim_response = youtube.search().list(
                    part="snippet",
                    q=similar_query,
                    type="video",
                    order="viewCount",
                    maxResults=10,
                    regionCode="KR",
                    relevanceLanguage="ko",
                    publishedAfter=datetime.now().strftime("%Y-%m-01T00:00:00Z"),
                ).execute()
                sim_video_ids = [item["id"]["videoId"] for item in sim_response.get("items", []) if item["id"].get("videoId")]
                registered_channel_ids = set(channel_ids)
                if sim_video_ids:
                    sim_detail = youtube.videos().list(part="snippet,statistics,contentDetails", id=",".join(sim_video_ids[:10])).execute()
                    for item in sim_detail.get("items", []):
                        if not _is_valid_video(item):
                            continue
                        snippet = item["snippet"]
                        if snippet.get("channelId") in registered_channel_ids:
                            continue
                        stats = item.get("statistics", {})
                        vid = item["id"]
                        all_videos.append({
                            "videoId": vid,
                            "title": snippet.get("title", ""),
                            "channelName": snippet.get("channelTitle", ""),
                            "thumbnail": f"https://img.youtube.com/vi/{vid}/hqdefault.jpg",
                            "viewCount": int(stats.get("viewCount", 0)),
                            "publishedAt": snippet.get("publishedAt", "")[:10],
                            "url": f"https://www.youtube.com/watch?v={vid}",
                            "source": "similar",
                            "sourceLabel": similar_query,
                        })

        seen = set()
        unique = []
        for v in all_videos:
            if v["videoId"] not in seen:
                seen.add(v["videoId"])
                unique.append(v)
        unique.sort(key=lambda x: x["viewCount"], reverse=True)
        return jsonify({"videos": unique[:5]})

    except Exception as e:
        print(f"맞춤 추천 오류: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/channel-info", methods=["GET"])
def channel_info():
    """영상 URL에서 채널 정보 추출"""
    api_key = os.getenv("YOUTUBE_API_KEY")
    video_url = request.args.get("url", "")

    import re
    match = re.search(r"(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})", video_url)
    if not match:
        return jsonify({"error": "올바른 유튜브 URL이 아닙니다"}), 400

    video_id = match.group(1)

    try:
        from googleapiclient.discovery import build
        youtube = build("youtube", "v3", developerKey=api_key)
        response = youtube.videos().list(part="snippet", id=video_id).execute()
        items = response.get("items", [])
        if not items:
            return jsonify({"error": "영상을 찾을 수 없습니다"}), 404

        snippet = items[0]["snippet"]
        return jsonify({
            "channelId": snippet["channelId"],
            "channelName": snippet["channelTitle"],
            "videoTitle": snippet["title"],
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/health", methods=["GET"])
def health():
    return jsonify({
        "status": "ok",
        "ai_connected": anthropic_client is not None,
        "db_connected": supabase_client is not None,
        "timestamp": datetime.now().isoformat(),
    })


def _do_analyze(analysis_id: str, video_urls: list):
    """백그라운드에서 영상 분석 수행"""
    try:
        videos = []
        for url in video_urls:
            video_id = extract_video_id(url)
            if not video_id:
                continue
            info = get_video_info(video_id)
            if not info:
                continue
            transcript = get_transcript(video_id)
            videos.append({
                **info,
                "url": url,
                "transcript": transcript,
                "viewRatio": 1.0,
            })

        if not videos:
            db_update_project(analysis_id, {
                "status": "error",
                "error_message": "유효한 영상을 찾을 수 없습니다",
            })
            return

        videos = calc_view_ratios(videos)
        topics = analyze_with_ai(videos)

        db_update_project(analysis_id, {
            "status": "topics_ready",
            "videos": videos,
            "topics": topics,
        })
        print(f"[분석 완료] {analysis_id} - {len(videos)}개 영상, {len(topics)}개 주제")
    except Exception as e:
        print(f"[분석 오류] {analysis_id}: {e}")
        db_update_project(analysis_id, {
            "status": "error",
            "error_message": str(e),
        })


@app.route("/api/analyze", methods=["POST"])
def analyze():
    """영상 분석 API (비동기)"""
    data = request.get_json()
    video_urls = data.get("video_urls", [])

    if len(video_urls) < 3:
        return jsonify({"error": "최소 3개의 영상 URL이 필요합니다"}), 400

    analysis_id = str(uuid.uuid4())
    project = {
        "id": analysis_id,
        "created_at": datetime.now().isoformat(),
        "video_urls": video_urls,
        "videos": [],
        "topics": [],
        "status": "analyzing",
    }
    db_save_project(project)

    thread = threading.Thread(target=_do_analyze, args=(analysis_id, video_urls))
    thread.daemon = True
    thread.start()

    return jsonify({
        "analysis_id": analysis_id,
        "status": "analyzing",
    })


@app.route("/api/retry-topics", methods=["POST"])
def retry_topics():
    """주제 다시 추천받기 API"""
    data = request.get_json()
    analysis_id = data.get("analysis_id")
    previous_topics = data.get("previous_topics", [])

    project = db_get_project(analysis_id)
    if not project:
        return jsonify({"error": "프로젝트를 찾을 수 없습니다"}), 404

    exclude_topics = [t.get("topic", "") for t in previous_topics if t.get("topic")]

    new_topics = analyze_with_ai(project["videos"], exclude_topics=exclude_topics)

    db_update_project(analysis_id, {
        "topics": new_topics,
        "status": "topics_ready",
    })

    return jsonify({
        "analysis_id": analysis_id,
        "topics": new_topics,
    })


def _do_generate(analysis_id: str, selected_topic: dict, videos: list, selected_topic_id: str):
    """백그라운드에서 콘텐츠 생성 수행"""
    try:
        result = generate_with_ai(selected_topic, videos)

        generated = {
            "id": str(uuid.uuid4()),
            "analysisId": analysis_id,
            "selectedTopic": selected_topic,
            **result,
        }

        db_update_project(analysis_id, {
            "status": "complete",
            "selected_topic_id": selected_topic_id,
            "generated_content": generated,
        })
        print(f"[생성 완료] {analysis_id}")
    except Exception as e:
        print(f"[생성 오류] {analysis_id}: {e}")
        db_update_project(analysis_id, {
            "status": "error",
            "error_message": str(e),
        })


@app.route("/api/generate", methods=["POST"])
def generate():
    """콘텐츠 생성 API (비동기)"""
    data = request.get_json()
    analysis_id = data.get("analysis_id")
    selected_topic_id = data.get("selected_topic_id")

    project = db_get_project(analysis_id)
    if not project:
        return jsonify({"error": "프로젝트를 찾을 수 없습니다"}), 404

    selected_topic = None
    for topic in project["topics"]:
        if topic["id"] == selected_topic_id:
            selected_topic = topic
            break

    if not selected_topic:
        return jsonify({"error": "선택된 주제를 찾을 수 없습니다"}), 404

    db_update_project(analysis_id, {"status": "generating"})

    thread = threading.Thread(
        target=_do_generate,
        args=(analysis_id, selected_topic, project["videos"], selected_topic_id),
    )
    thread.daemon = True
    thread.start()

    return jsonify({"status": "generating", "analysis_id": analysis_id})


# ============================================================
# 제작 API 엔드포인트 (Steps 5-10)
# ============================================================

@app.route("/api/confirm", methods=["POST"])
def confirm_content():
    """콘텐츠 확정 API (Step 5)"""
    data = request.get_json()
    analysis_id = data.get("analysis_id")
    confirmed = data.get("confirmed_content")

    if not analysis_id or not confirmed:
        return jsonify({"error": "analysis_id와 confirmed_content 필요"}), 400

    project = db_get_project(analysis_id)
    if not project:
        return jsonify({"error": "프로젝트를 찾을 수 없습니다"}), 404

    db_update_project(analysis_id, {
        "status": "confirmed",
        "confirmed_content": confirmed,
    })

    return jsonify({"status": "confirmed", "analysis_id": analysis_id})


def _do_produce(analysis_id: str, confirmed: dict):
    """백그라운드에서 전체 제작 수행 (썸네일 → PPT → TTS → 영상)"""
    assets = {
        "thumbnailUrl": None,
        "pptUrl": None,
        "ttsUrl": None,
        "videoUrl": None,
    }

    try:
        # 1. 썸네일 생성
        print(f"[제작] {analysis_id} - 썸네일 생성 중...")
        thumb_file = generate_thumbnail_image(
            confirmed.get("thumbnailText", "썸네일"),
            analysis_id
        )
        assets["thumbnailUrl"] = f"/api/assets/{thumb_file}"
        db_update_project(analysis_id, {"production_assets": assets})

        # 2. PPT 생성
        print(f"[제작] {analysis_id} - PPT 생성 중...")
        script = confirmed.get("script", {})
        ppt_file = generate_ppt(
            confirmed.get("title", "제목 없음"),
            script,
            analysis_id
        )
        assets["pptUrl"] = f"/api/assets/{ppt_file}"
        db_update_project(analysis_id, {"production_assets": assets})

        # 3. TTS 생성
        print(f"[제작] {analysis_id} - TTS 생성 중...")
        tts_file = generate_tts(script, analysis_id)
        assets["ttsUrl"] = f"/api/assets/{tts_file}"
        db_update_project(analysis_id, {"production_assets": assets})

        # 4. 영상 합성
        print(f"[제작] {analysis_id} - 영상 합성 중...")
        try:
            video_file = generate_video(analysis_id, ppt_file, tts_file)
            assets["videoUrl"] = f"/api/assets/{video_file}"
        except Exception as ve:
            print(f"[영상 합성 경고] {ve} - FFmpeg 미설치일 수 있음")
            assets["videoUrl"] = None

        db_update_project(analysis_id, {
            "status": "production_done",
            "production_assets": assets,
        })
        print(f"[제작 완료] {analysis_id}")

    except Exception as e:
        print(f"[제작 오류] {analysis_id}: {e}")
        # 부분 에셋이라도 저장
        db_update_project(analysis_id, {
            "status": "production_done",
            "production_assets": assets,
            "error_message": str(e),
        })


@app.route("/api/produce", methods=["POST"])
def produce():
    """제작 시작 API (Step 6-9, 비동기)"""
    data = request.get_json()
    analysis_id = data.get("analysis_id")

    project = db_get_project(analysis_id)
    if not project:
        return jsonify({"error": "프로젝트를 찾을 수 없습니다"}), 404

    confirmed = project.get("confirmed_content")
    if not confirmed:
        return jsonify({"error": "확정된 콘텐츠가 없습니다. 먼저 /api/confirm을 호출하세요"}), 400

    db_update_project(analysis_id, {
        "status": "producing",
        "production_assets": {
            "thumbnailUrl": None,
            "pptUrl": None,
            "ttsUrl": None,
            "videoUrl": None,
        },
    })

    thread = threading.Thread(target=_do_produce, args=(analysis_id, confirmed))
    thread.daemon = True
    thread.start()

    return jsonify({"status": "producing", "analysis_id": analysis_id})


@app.route("/api/assets/<filename>", methods=["GET"])
def serve_asset(filename):
    """제작된 에셋 파일 서빙"""
    # 경로 조작 방지
    safe_name = os.path.basename(filename)
    filepath = ASSETS_DIR / safe_name

    if not filepath.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다"}), 404

    # MP4 영상은 명시적 MIME type 지정 (브라우저 재생 호환성)
    mimetype = None
    if safe_name.endswith(".mp4"):
        mimetype = "video/mp4"
    elif safe_name.endswith(".mp3"):
        mimetype = "audio/mpeg"
    elif safe_name.endswith(".pptx"):
        mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    return send_file(str(filepath), mimetype=mimetype)


@app.route("/api/upload-youtube", methods=["POST"])
def upload_youtube():
    """YouTube 업로드 API (Step 10)

    참고: 실제 YouTube 업로드는 OAuth 2.0 인증이 필요합니다.
    현재는 업로드 준비 상태만 반환합니다.
    """
    data = request.get_json()
    analysis_id = data.get("analysis_id")

    project = db_get_project(analysis_id)
    if not project:
        return jsonify({"error": "프로젝트를 찾을 수 없습니다"}), 404

    assets = project.get("production_assets", {})
    confirmed = project.get("confirmed_content", {})

    if not assets or not assets.get("videoUrl"):
        return jsonify({"error": "업로드할 영상이 없습니다. 먼저 제작을 완료해주세요"}), 400

    # YouTube OAuth 토큰 확인
    # 현재는 OAuth가 구현되지 않았으므로 안내 메시지 반환
    return jsonify({
        "status": "ready",
        "message": "YouTube 업로드를 위해서는 Google OAuth 인증이 필요합니다. 설정 페이지에서 Google 계정을 연동해주세요.",
        "upload_data": {
            "title": confirmed.get("title", ""),
            "description": f"주제: {confirmed.get('topic', {}).get('topic', '')}\n\n이 영상은 AI 도구로 제작되었습니다.",
            "video_file": assets.get("videoUrl", ""),
            "thumbnail_file": assets.get("thumbnailUrl", ""),
        },
    })


@app.route("/api/project/<project_id>", methods=["GET"])
def get_project(project_id):
    """프로젝트 상태 조회"""
    project = db_get_project(project_id)
    if not project:
        return jsonify({"error": "프로젝트를 찾을 수 없습니다"}), 404
    return jsonify(project)


# ============================================================
# 실행
# ============================================================

if __name__ == "__main__":
    port = int(os.getenv("PORT", 5000))
    print(f"서버 시작: http://0.0.0.0:{port}")
    print(f"  AI 연동: {'활성' if anthropic_client else '비활성 (ANTHROPIC_API_KEY 필요)'}")
    print(f"  DB 연동: {'활성' if supabase_client else '비활성 (인메모리 모드)'}")
    app.run(host="0.0.0.0", port=port, debug=True)
